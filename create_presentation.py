import pptx
from pptx import Presentation

# This script generates a basic PowerPoint deck for the RAG project.
# Install dependency: pip install python-pptx

slides = [
    {"title": "Projet RAG - API Culture Paris", "bullets": [
        "Objectif : rechercher des événements culturels parisiens via un assistant AI",
        "Contexte : utilisation de données OpenData et de modèles Mistral"
    ]},
    {"title": "Architecture technique", "bullets": [
        "Collecte & nettoyage des données OpenDataSoft",
        "Indexation vectorielle FAISS avec mistral-embed",
        "API REST FastAPI exposant /ask et /rebuild",
        "Conteneurisation Docker (image exécutable)"
    ]},
    {"title": "Données et modèles", "bullets": [
        "Source : événements publics OpenAgenda (Paris)",
        "Embeddings : Mistral-Embed",
        "LLM de génération : Mistral-Small-Latest",
        "Stockage : index FAISS local"
    ]},
    {"title": "Tests et évaluation", "bullets": [
        "8 tests pytest couvrant comportement API et logique RAG",
        "Scénarios d'usage pour la démonstration",
        "Évaluation manuelle via script de test et logs" 
    ]},
    {"title": "Démonstration en direct", "bullets": [
        "1. Construction de l'index (build_index.sh)",
        "2. Lancement de l'API (uvicorn ou conteneur)",
        "3. Questions exemples : jazz, gratuit, hors-contexte",
        "4. Réponse retournée en quelques centièmes de seconde"
    ]},
    {"title": "Perspectives d'amélioration", "bullets": [
        "Prise en compte de données hors-ligne / planifiées",
        "Optimisation du chunking et du rappel de contexte",
        "Cache pour réponses fréquentes, mise à l'échelle Docker",
        "Interface utilisateur simple ou Chatbot front-end"
    ]},
    {"title": "Questions techniques & métier", "bullets": [
        "Pourquoi un RAG ? Avantages vs chatbot classique",
        "Choix de FastAPI et Docker pour la portabilité",
        "Sécurité des clés API et gestion des dépendances",
        "Plan de continuité pour la démo hors connexion"
    ]}
]

prs = Presentation()
for slide in slides:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = slide["title"]
    body = s.shapes.placeholders[1].text_frame
    for bullet in slide["bullets"]:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 0

prs.save("Projet_RAG.pptx")
print("Presentation created: Projet_RAG.pptx")
